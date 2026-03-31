#!/usr/bin/env python3
"""
Genera la presentación PowerPoint para el POC NuGet POS.
Ejecutar: python3 docs/generate-pptx.py
Salida: docs/POC-NuGet-POS-Presentacion.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors ──────────────────────────────────────────────────────────────────
SIESA_BLUE = RGBColor(0x1A, 0x23, 0x7E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
YELLOW_LIGHT = RGBColor(0xFF, 0xF9, 0xC4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide():
    """Add blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_title_bar(slide, text, subtitle=None):
    """Add a colored title bar at top."""
    # Bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SIESA_BLUE
    shape.line.fill.background()
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)


def add_body_text(slide, text, left=0.5, top=1.5, width=12, size=18, bold=False, color=DARK_GRAY):
    """Add body text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(8)
    return txBox


def add_bullet_list(slide, items, left=0.5, top=1.5, width=12, size=18):
    """Add bulleted list."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    return txBox


def add_table(slide, headers, rows, left=0.5, top=1.5, width=12, row_height=0.45):
    """Add a styled table."""
    cols = len(headers)
    num_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(
        num_rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(row_height * num_rows))
    tbl = tbl_shape.table
    # Set column widths evenly
    col_w = Emu(int(Inches(width) / cols))
    for c in range(cols):
        tbl.columns[c].width = col_w
    # Header row
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SIESA_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
    return tbl


def add_box(slide, text, left, top, width, height, bg_color, text_color=DARK_GRAY, size=14):
    """Add colored box with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = text_color
        if i == 0:
            p.font.bold = True


def add_arrow_text(slide, text, left, top, size=24):
    """Add arrow text between elements."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(1.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = MED_GRAY


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Portada
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
# Full blue background
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = SIESA_BLUE
bg.line.fill.background()
add_body_text(s, "POC", 0.5, 1.5, 12, 60, True, WHITE)
add_body_text(s, "Distribución de Módulos POS\nvia NuGet", 0.5, 2.5, 12, 36, False, RGBColor(0xBB, 0xDE, 0xFB))
add_body_text(s, "Un solo código → múltiples despliegues", 0.5, 4.0, 12, 24, False, RGBColor(0x90, 0xCA, 0xF9))
add_body_text(s, "SIESA  •  Equipo Técnico  •  Marzo 2026", 0.5, 5.5, 12, 18, False, RGBColor(0x64, 0xB5, 0xF6))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: El Problema
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "El Problema")
add_bullet_list(s, [
    "El POS necesita correr en tiendas físicas (offline) Y en la nube (centralizado)",
    "Hoy: código duplicado entre proyectos para cada contexto",
    "Cambios se aplican manualmente → inconsistencias y bugs",
    "Sin estrategia de distribución: cada host es un proyecto independiente",
    "Esfuerzo duplicado de desarrollo y mantenimiento",
], top=1.5, size=20)
add_box(s, "RESULTADO\nInconsistencias, bugs,\nesfuerzo duplicado", 9, 4.5, 3.5, 1.8, RGBColor(0xFF, 0xEB, 0xEE), RED, 16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: La Propuesta
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "La Propuesta", "Un solo repositorio, múltiples despliegues via NuGet")
# Flow boxes
add_box(s, "Código Fuente\n(Módulos)", 0.5, 2, 3, 1.5, RGBColor(0xF3, 0xE5, 0xF5))
add_arrow_text(s, "→", 3.6, 2.3)
add_box(s, "dotnet pack\n8 paquetes NuGet", 4.5, 2, 3, 1.5, YELLOW_LIGHT)
add_arrow_text(s, "→", 7.6, 2.3)
add_box(s, "NuGet Feed\n(Local / GitHub)", 8.5, 2, 3, 1.5, RGBColor(0xFF, 0xE0, 0xB2))
# Deployments
add_box(s, "CloudHub (Nube)\nTodos los módulos\nSource of truth", 2, 4.5, 4, 1.8, BLUE_LIGHT)
add_box(s, "LocalPOS (Tienda)\nSubset de módulos\nOffline-first", 7, 4.5, 4, 1.8, GREEN_LIGHT)
add_arrow_text(s, "↓", 9.5, 3.7)
add_arrow_text(s, "↓", 3.5, 3.7)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Arquitectura General
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Arquitectura General", "CloudHub (nube central) ↔ LocalPOS (tienda)")
# Cloud
add_box(s, "☁️ CLOUD HUB (:5200)\nOrders + Products + Payments\n+ Sync Receivers\n+ (Reportes, Admin, Analytics...)\nDB: pos_cloud", 0.5, 1.5, 5.5, 2.5, BLUE_LIGHT, DARK_GRAY, 16)
# Internet
add_body_text(s, "🌐  INTERNET (puede caerse en cualquier momento)", 0.5, 4.2, 12, 14, False, ORANGE)
# Local
add_box(s, "🏪 LOCAL POS (:5100)\nOrders + Products + Payments\n+ Sync Worker (Outbox)\nDB: pos_local", 0.5, 4.8, 5.5, 2, GREEN_LIGHT, DARK_GRAY, 16)
# Key points
add_bullet_list(s, [
    "CloudHub: sistema central con TODOS los módulos",
    "LocalPOS: POS de tienda con módulos necesarios",
    "Sync: local → cloud via Outbox Pattern (cada 10s)",
    "Ambos comparten el MISMO código (NuGet packages)",
    "En producción: CloudHub tendría más módulos",
], left=6.5, top=1.5, width=6, size=16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Módulos NuGet
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Módulos NuGet", "8 paquetes distribuibles independientemente")
add_table(s,
    ["Paquete", "Tipo", "Contenido"],
    [
        ["Pos.SharedKernel", "Shared", "BaseEntity, ULID, IEventBus, SyncOutbox"],
        ["Pos.Infrastructure.Postgres", "Infrastructure", "MigrationRunner, EF Core config"],
        ["Pos.Orders.Contracts", "Contracts", "OrderDto, IOrderRepository, requests"],
        ["Pos.Orders.Core", "Core", "Order entity, EfOrderRepository, endpoints"],
        ["Pos.Products.Contracts", "Contracts", "ProductDto, CategoryDto, interfaces"],
        ["Pos.Products.Core", "Core", "Product/Category entities, repos, endpoints"],
        ["Pos.Payments.Contracts", "Contracts", "PaymentDto, IPaymentRepository"],
        ["Pos.Payments.Core", "Core", "Payment entity, EfPaymentRepository, endpoints"],
    ],
    top=1.5, row_height=0.5)
add_body_text(s, "Cada módulo se divide en Contracts (interfaces + DTOs) y Core (implementación).\nLos consumidores pueden depender solo de Contracts para bajo acoplamiento.", 0.5, 6.2, 12, 14, False, MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Dual-Mode
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "NuGet Dual-Mode", "El mecanismo central del POC")
add_body_text(s, 'Directory.Build.props controla cómo se resuelven las dependencias:', 0.5, 1.5, 12, 18)
add_table(s,
    ["Modo", "Flag", "Dependencias", "Uso"],
    [
        ["Desarrollo", "UseProjectReference=true", "ProjectReference", "Debug directo en módulos, F12 al código"],
        ["CI/CD", "UseProjectReference=false", "PackageReference (NuGet)", "Distribución via paquetes publicados"],
    ],
    top=2.3, row_height=0.6)
add_body_text(s, "Cómo funciona:", 0.5, 3.8, 12, 18, True)
add_bullet_list(s, [
    "Cada .csproj tiene un bloque <Choose> con ambas opciones (ProjectRef y PackageRef)",
    "Directory.Build.props define el default: UseProjectReference=true",
    "En CI/CD se pasa -p:UseProjectReference=false para forzar modo NuGet",
    "Los hosts NO saben si consumen código fuente o paquetes — es transparente",
    "PosModuleVersion centraliza la versión para que pack y restore coincidan",
], top=4.2, size=16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Sync Outbox Pattern
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Sincronización: Outbox Pattern", "Local-first + consistencia eventual")
# Flow 1
add_box(s, "FLUJO 1: Write Path (<1ms)\nPOS → API → Repository → TRANSACCIÓN:\nINSERT dominio + INSERT outbox → 200 OK", 0.5, 1.5, 5.5, 1.5, RGBColor(0xFF, 0xF3, 0xE0), DARK_GRAY, 14)
# Flow 2
add_box(s, "FLUJO 2: Sync Push (cada 10s)\nTimer → ¿Internet? → SELECT outbox\nWHERE synced_at IS NULL LIMIT 50\n→ POST CloudHub → UPDATE synced_at", 0.5, 3.3, 5.5, 1.7, RGBColor(0xE8, 0xEA, 0xF6), DARK_GRAY, 14)
# Flow 3
add_box(s, "FLUJO 3: Pull Path (futuro)\nCloud → GET /changes?since=watermark\n→ Local UPSERT\n⚠️ No implementado en POC", 0.5, 5.3, 5.5, 1.5, RGBColor(0xFC, 0xE4, 0xEC), DARK_GRAY, 14)
# Key concepts
add_body_text(s, "Conceptos clave:", 6.5, 1.5, 6, 18, True)
add_bullet_list(s, [
    "Local-first: la operación local completa en <1ms",
    "Idempotencia: UPSERT por ULID primary key",
    "Tolerancia a fallas: si cloud cae, outbox retiene",
    "Retry automático en el siguiente ciclo",
    "Sin infraestructura adicional (usa la misma DB)",
    "Datos Local→Cloud: órdenes, productos, pagos",
    "Datos Cloud→Local: catálogos, precios (futuro)",
], left=6.5, top=2.2, width=6, size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Demo Paso 1
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Demo: Paso 1", "Levantar el entorno con NuGet real")
add_body_text(s, "1. Levantar PostgreSQL:", 0.5, 1.5, 12, 18, True)
add_box(s, "docker-compose up -d", 0.5, 2.0, 6, 0.6, RGBColor(0x26, 0x32, 0x38), WHITE, 16)
add_body_text(s, "2. Empaquetar módulos como NuGet:", 0.5, 2.9, 12, 18, True)
add_box(s, "dotnet pack -c Release\nls artifacts/nupkg/   # → 8 paquetes .nupkg", 0.5, 3.4, 8, 0.9, RGBColor(0x26, 0x32, 0x38), WHITE, 16)
add_body_text(s, "3. Arrancar hosts consumiendo NuGet (NO código fuente):", 0.5, 4.6, 12, 18, True)
add_box(s, "# Terminal 1 (Cloud Hub)\ndotnet run --project src/Hosts/Pos.Host.CloudHub -c Release -p:UseProjectReference=false\n\n# Terminal 2 (Local POS)\ndotnet run --project src/Hosts/Pos.Host.LocalPOS -c Release -p:UseProjectReference=false", 0.5, 5.1, 11, 1.8, RGBColor(0x26, 0x32, 0x38), WHITE, 14)
add_box(s, "El flag -p:UseProjectReference=false\nhace que los hosts consuman\nlos .nupkg del feed local", 9.5, 1.5, 3.3, 1.5, YELLOW_LIGHT, DARK_GRAY, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Demo Paso 2
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Demo: Paso 2", "CRUD en POS local + Sync automático")
add_body_text(s, "Crear datos en LocalPOS (localhost:5100):", 0.5, 1.5, 12, 18, True)
add_bullet_list(s, [
    "POST /api/categories → Crear categoría 'Bebidas'",
    "POST /api/products → Crear producto 'Café Americano' ($5,500)",
    "POST /api/orders → Crear orden con líneas de productos",
    "POST /api/orders/{id}/close → Cerrar la orden",
    "POST /api/payments → Registrar pago en efectivo",
], top=2.2, size=16)
add_body_text(s, "Esperar ~10 segundos y verificar en CloudHub (localhost:5200):", 0.5, 4.5, 12, 18, True)
add_box(s, "curl http://localhost:5200/api/products   # ¡Producto aparece!\ncurl http://localhost:5200/api/orders     # ¡Orden aparece!", 0.5, 5.1, 10, 0.9, RGBColor(0x26, 0x32, 0x38), WHITE, 16)
add_box(s, "Los datos creados\nen la tienda (:5100)\naparecen en la nube (:5200)\nautomáticamente", 9, 4.0, 3.8, 1.5, GREEN_LIGHT, GREEN, 14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Demo Paso 3 — Cambio en vivo
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Demo: Paso 3 — Cambio en Vivo", "Modificar módulo → re-pack → ambos hosts actualizados")
add_body_text(s, "1. Detener ambos hosts", 0.5, 1.4, 12, 16, True)
add_body_text(s, "2. Agregar campo 'PriceFormatted' al ProductDto (1 archivo) + ToDto (1 archivo)", 0.5, 1.8, 12, 16, True)
add_body_text(s, "3. Re-empaquetar y re-arrancar:", 0.5, 2.3, 12, 16, True)
add_box(s, "dotnet pack -c Release\ndotnet run --project src/Hosts/Pos.Host.CloudHub -c Release -p:UseProjectReference=false\ndotnet run --project src/Hosts/Pos.Host.LocalPOS -c Release -p:UseProjectReference=false", 0.5, 2.8, 10, 1.3, RGBColor(0x26, 0x32, 0x38), WHITE, 14)
add_body_text(s, "4. Verificar:", 0.5, 4.3, 12, 16, True)
add_box(s, 'curl localhost:5100/api/products | jq \'.[0].priceFormatted\'   # → "$5,500 COP"\ncurl localhost:5200/api/products | jq \'.[0].priceFormatted\'   # → "$5,500 COP"', 0.5, 4.8, 10, 0.9, RGBColor(0x26, 0x32, 0x38), WHITE, 14)
add_box(s, "Se modificó el código\nen UN solo lugar.\nSe re-empaquetó como NuGet.\nAMBOS hosts muestran\nel nuevo campo.\n\nEsto es exactamente lo que\nharía CI/CD en producción.", 8.5, 1.5, 4.3, 3, YELLOW_LIGHT, DARK_GRAY, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Demo Paso 4 — Resiliencia
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Demo: Paso 4 — Resiliencia Offline", "¿Qué pasa si la nube se cae?")
add_bullet_list(s, [
    "1. Detener CloudHub (Ctrl+C)",
    "2. Crear un producto nuevo en LocalPOS → funciona normalmente",
    "3. Logs muestran: 'Cloud unreachable... will retry next cycle'",
    "4. Re-arrancar CloudHub",
    "5. En ~10 segundos el producto se sincroniza automáticamente",
], top=1.5, size=20)
add_box(s, "CONCLUSIÓN\n\nSi la nube se cae:\n• El POS sigue funcionando\n• Los datos NO se pierden\n• El outbox retiene todo\n• Sync automático al reconectar", 7.5, 3.5, 5, 3, GREEN_LIGHT, GREEN, 16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Pros y Contras
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Pros y Contras")
add_table(s,
    ["", "Pros ✅", "Contras ⚠️"],
    [
        ["Código", "Una sola base de código para todos los contextos", "Requiere CI/CD configurado para distribución"],
        ["Desarrollo", "Debug directo con ProjectReference (F12)", "Developers deben entender el mecanismo dual"],
        ["Distribución", "Versionado semántico via NuGet nativo", "Version drift si hosts no se actualizan juntos"],
        ["Sync", "Local-first, tolerante a fallas, sin infra extra", "Polling (10s latencia), sin DLQ, sin retry avanzado"],
        ["Módulos", "Desacoplados, cada uno con su schema DB", "Sin FK cross-schema (by design)"],
        ["Evolución", "Path gradual a microservicios", "Requiere disciplina en las interfaces (Contracts)"],
    ],
    top=1.4, row_height=0.65)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Riesgos y Mitigaciones
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Riesgos y Mitigaciones")
add_table(s,
    ["Riesgo", "Prob.", "Impacto", "Mitigación"],
    [
        ["Version drift entre módulos NuGet", "Alta", "Alto", "Versión centralizada en Directory.Build.props + CI valida"],
        ["Pérdida de datos en sync", "Media", "Alto", "Marcar synced DESPUÉS de ACK + upsert idempotente"],
        ["Sin autenticación en endpoints", "Alta (POC)", "Alto (Prod)", "API key como primer paso post-MVP"],
        ["Sync entries acumuladas (sin DLQ)", "Media", "Medio", "Monitoreo + DLQ post-MVP"],
        [".NET 10 / EF Core 10 estabilidad", "Media", "Medio", "Versiones fijas + rollback path a .NET 9"],
        ["Breaking changes en Contracts", "Media", "Alto", "Semantic versioning + CI que compila consumidores"],
    ],
    top=1.4, row_height=0.65)
add_body_text(s, "Ver docs/risk-matrix.md para análisis detallado con propuestas post-MVP", 0.5, 6.5, 12, 13, False, MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Decisiones Arquitectónicas
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Decisiones Arquitectónicas (ADRs)")
add_table(s,
    ["ADR", "Decisión", "Alternativas descartadas"],
    [
        ["001", "Modular Monolith + NuGet Distribution", "Microservicios, Git submodules, shared library"],
        ["002", "Outbox Pattern (polling) para sync", "CDC/Debezium, HTTP directo, Event Sourcing, RabbitMQ"],
        ["003", "Multi-Schema PostgreSQL por módulo", "DBs separadas, schema único, prefijos de tabla"],
    ],
    top=1.5, row_height=0.7)
add_body_text(s, "¿Por qué NuGet y no microservicios?", 0.5, 3.8, 12, 18, True)
add_bullet_list(s, [
    "Microservicios requieren: API Gateway, message broker, service mesh, service discovery",
    "Overhead de infraestructura desproporcionado para el tamaño actual del equipo",
    "Riesgo de 'distributed monolith' si no se hace bien",
    "NuGet ofrece distribución SIN el overhead → path gradual a microservicios cuando sea necesario",
], top=4.3, size=16)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15: Stack Técnico
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Stack Técnico")
add_table(s,
    ["Componente", "Tecnología", "Versión", "Rol"],
    [
        ["Runtime", ".NET", "10.0", "Framework principal"],
        ["ORM", "Entity Framework Core", "10.0", "Acceso a datos, multi-schema"],
        ["Base de datos", "PostgreSQL", "17", "Almacenamiento, schemas aislados"],
        ["IDs", "ULID", "1.3.4", "Identificadores ordenables por tiempo"],
        ["API Docs", "Scalar (OpenAPI)", "2.6", "UI interactiva para explorar APIs"],
        ["Paquetes", "NuGet", "—", "Distribución de módulos"],
        ["CI/CD", "GitHub Actions", "—", "Pack + Push + Validate"],
        ["Contenedores", "Docker Compose", "—", "PostgreSQL local para desarrollo"],
        ["Patrones", "Records inmutables (C#)", "—", "Sin mutación, thread-safe"],
    ],
    top=1.4, row_height=0.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16: Roadmap Post-MVP
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Roadmap Post-MVP")
# Phase boxes
add_box(s, "FASE 1: Hardening\n\n• Autenticación (JWT + API key)\n• Input validation (FluentValidation)\n• EF Core Migrations\n• Rate limiting", 0.5, 1.5, 3, 3, BLUE_LIGHT, DARK_GRAY, 13)
add_box(s, "FASE 2: Resiliencia\n\n• Retry con Polly + circuit breaker\n• Dead-letter queue para sync\n• Health checks avanzados\n• Outbox backlog monitoring", 3.7, 1.5, 3, 3, GREEN_LIGHT, DARK_GRAY, 13)
add_box(s, "FASE 3: Observabilidad\n\n• OpenTelemetry (tracing)\n• Métricas de sync\n• Dashboard de estado\n• Alertas", 6.9, 1.5, 3, 3, YELLOW_LIGHT, DARK_GRAY, 13)
add_box(s, "FASE 4: Escalabilidad\n\n• Multi-tenancy (por tienda)\n• Semantic versioning + ApiCompat\n• CDC (Debezium) para alto volumen\n• Pull path (Cloud→Local)", 10.1, 1.5, 3, 3, RGBColor(0xFC, 0xE4, 0xEC), DARK_GRAY, 13)
# Arrows
add_arrow_text(s, "→", 3.4, 2.5)
add_arrow_text(s, "→", 6.6, 2.5)
add_arrow_text(s, "→", 9.8, 2.5)
# Timeline
add_body_text(s, "Prioridad: Fase 1 es prerequisito para producción. Fases 2-4 son incrementales.", 0.5, 5, 12, 16, True, SIESA_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17: Conclusión
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_title_bar(s, "Conclusión y Recomendación")
add_body_text(s, "El POC demuestra que es viable:", 0.5, 1.5, 12, 20, True)
add_bullet_list(s, [
    "Mantener UNA SOLA base de código para múltiples contextos de despliegue",
    "Distribuir módulos como NuGet con switch automático dev/CI",
    "Sincronizar datos local → cloud con outbox pattern sin infraestructura adicional",
    "Operar offline-first con resiliencia automática ante caídas de red",
    "Evolucionar gradualmente de monolito modular a microservicios",
], top=2.2, size=18)
add_box(s, "RECOMENDACIÓN\n\nProceder con Fase 1 (Hardening)\ny adoptar este patrón como estándar\npara los módulos POS de SIESA", 3, 5, 7, 2, SIESA_BLUE, WHITE, 18)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "POC-NuGet-POS-Presentacion.pptx")
prs.save(output_path)
print(f"✅ Presentación generada: {output_path}")
print(f"   {len(prs.slides)} slides")
